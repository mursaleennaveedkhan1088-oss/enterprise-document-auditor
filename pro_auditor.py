import streamlit as st
from spellchecker import SpellChecker
import textstat
import pypdf
import docx
from pptx import Presentation
import language_tool_python
import re

# Set up professional widescreen layout
st.set_page_config(page_title="Enterprise Document Auditor", page_icon="🏢", layout="wide")

st.title("🏢 Enterprise Document & Quality Assurance Auditor")
st.markdown("A **digitalisation and workflow automation app** built with Python to streamline document reviews, scanning multi-format corporate and technical files for grammar, readability, and **formatting consistency**.")

# --- 1. CLOUD GRAMMAR ENGINE (No Local Java Required) ---
@st.cache_resource
def get_grammar_tool():
    try:
        return language_tool_python.LanguageToolPublicAPI('en-US')
    except Exception:
        return None

# --- 2. DOCUMENT EXTRACTION FUNCTIONS ---
def extract_from_pdf(file):
    reader = pypdf.PdfReader(file)
    text = ""
    for page in reader.pages:
        extracted = page.extract_text()
        if extracted:
            text += extracted + "\n"
    return text, [], []

def extract_from_docx(file):
    doc = docx.Document(file)
    text = ""
    fonts_found, sizes_found = set(), set()
    for p in doc.paragraphs:
        text += p.text + "\n"
        for run in p.runs:
            if run.font.name: fonts_found.add(run.font.name)
            if run.font.size: sizes_found.add(f"{run.font.size.pt} pt")
    return text, list(fonts_found), list(sizes_found)

def extract_from_pptx(file):
    prs = Presentation(file)
    text = ""
    fonts_found, sizes_found = set(), set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    text += p.text + "\n"
                    for run in p.runs:
                        if run.font.name: fonts_found.add(run.font.name)
                        if run.font.size: sizes_found.add(f"{run.font.size.pt} pt")
    return text, list(fonts_found), list(sizes_found)

# --- 3. DRAG & DROP FILE UPLOADER (Text Pasting Removed) ---
uploaded_file = st.file_uploader(
    "📁 Drag and drop your Word (.docx), PDF, or PowerPoint (.pptx) file here:",
    type=["txt", "pdf", "docx", "pptx"],
    help="Supports Word reports, PDF contracts/summaries, PowerPoint presentations, and plain text."
)

# --- 4. CORPORATE & TECHNICAL TERMINOLOGY WHITELIST ---
with st.expander("⚙️ Dictionary Whitelist (Ignored Technical, Chemical & Name Terms)"):
    st.write("These industry terms, acronyms, names, and operational variables will be ignored by the spellchecker:")
    # Pre-loaded with engineering, chemical adsorption, corporate terms, and common names
    default_whitelist = (
        "kpi, roi, saas, workflow, automation, dashboard, analytics, pptx, docx, pdf, "
        "online, platform, portal, crm, erp, telemetry, scada, q1, q2, q3, q4, yoy, "
        "stakeholders, roadmap, onboarding, scalable, infrastructure, frontend, backend, "
        "arduino, zeolites, zeolite, psa, o2, lpm, bar, labview, simscape, haptics, matlab, "
        "autocad, microcontrollers, coursera, www, uno, adsorption, adsorber, adsorben, adsorbent, "
        "concentrator, valves, solenoid, fyp, thesis, oom202, aluminosilicate, airtac, aiche, acs, "
        "adc, addr, ackley, alue, analysing, analyzing, abdullah, ajmal, amarullah, muzammil, riaz, rafique, humayun"
    )
    user_whitelist_input = st.text_input("Custom Ignore List (separate by commas):", value=default_whitelist)
    custom_whitelist = {w.strip().lower() for w in user_whitelist_input.split(",") if w.strip()}

# Determine text source from file only
document_text = ""
detected_fonts, detected_sizes = [], []
file_type = "None"

if uploaded_file is not None:
    file_ext = uploaded_file.name.split(".")[-1].lower()
    file_type = file_ext.upper()
    if file_ext == "txt":
        document_text = uploaded_file.read().decode("utf-8", errors="ignore")
    elif file_ext == "pdf":
        document_text, detected_fonts, detected_sizes = extract_from_pdf(uploaded_file)
    elif file_ext == "docx":
        document_text, detected_fonts, detected_sizes = extract_from_docx(uploaded_file)
    elif file_ext == "pptx":
        document_text, detected_fonts, detected_sizes = extract_from_pptx(uploaded_file)

# --- 5. DEEP DOCUMENT AUDIT ---
if st.button("🔍 Run Quality Assurance Audit", type="primary"):
    if uploaded_file is None or not document_text.strip():
        st.warning("⚠️ Please upload a valid document file above before running the audit!")
    else:
        st.divider()
        tab1, tab2, tab3 = st.tabs(["✍️ Grammar & Tone", "🎨 Brand Formatting Consistency", "📊 Readability & Workflow Stats"])
        
        # --- TAB 1: GRAMMAR & SPELLING ---
        with tab1:
            col_gram, col_spell = st.columns(2)
            
            with col_gram:
                st.subheader("📝 Grammar, Punctuation & Professional Tone")
                tool = get_grammar_tool()
                if tool:
                    with st.spinner("Connecting to Cloud NLP Engine (scanning first 15,000 characters to prevent cloud timeout)..."):
                        try:
                            # Cap text to 15,000 characters to prevent HTTP 413/500 cloud API crash on massive PDF reports
                            safe_text_sample = document_text[:15000]
                            matches = tool.check(safe_text_sample)
                            
                            if len(document_text) > 15000:
                                st.info("ℹ️ Note: Due to cloud API rate limits on massive files, grammar check was applied to the first ~2,500 words.")
                            
                            if not matches:
                                st.success("✔ Perfect! No grammatical, phrasing, or punctuation defects found in the scanned section.")
                            else:
                                st.error(f"Found {len(matches)} grammar / style issue(s):")
                                for error in matches[:10]:
                                    st.markdown(f"**Issue:** {error.message}")
                                    st.markdown(f"• *Context:* \"...`{error.context}`...\"")
                                    if error.replacements:
                                        st.markdown(f"• *Best Fix:* **`{error.replacements[0]}`**")
                                    st.divider()
                        except Exception as e:
                            st.warning("⚠️ Cloud Grammar API timed out or rejected the request due to network limits. Please try running the check again or upload a smaller chapter.")
                else:
                    st.error("Could not connect to Cloud Grammar API. Please check your internet connection.")
            
            with col_spell:
                st.subheader("❌ Word Typo Check")
                spell = SpellChecker()
                words_list = spell.split_words(document_text)
                
                # Ignore numbers, alphanumeric codes, short words <= 2 letters, and custom whitelist terms
                clean_words = [w for w in words_list if not re.search(r'\d', w) and len(w) > 2]
                
                raw_misspelled = spell.unknown(clean_words)
                true_misspelled = {word for word in raw_misspelled if word.lower() not in custom_whitelist}
                
                if not true_misspelled:
                    st.success("✔ No typos found! (Technical terms, names & operational codes ignored)")
                else:
                    st.warning(f"Found {len(true_misspelled)} typo(s):")
                    for word in sorted(list(true_misspelled))[:15]:
                        st.markdown(f"• **{word}** → *Did you mean: `{spell.correction(word)}`?*")
        
        # --- TAB 2: BRAND FORMATTING CONSISTENCY ---
        with tab2:
            st.subheader(f"🎨 Brand & Document Formatting Audit ({file_type})")
            st.write("Catches copy-paste formatting anomalies and irregular text sizing across paragraphs:")
            
            if file_type in ["DOCX", "PPTX"]:
                c1, c2 = st.columns(2)
                with c1:
                    st.write("**Detected Font Families:**")
                    if detected_fonts:
                        for f in detected_fonts: st.code(f)
                        if len(detected_fonts) > 2: 
                            st.warning("⚠️ High Font Variation: Corporate guidelines usually require sticking to 1 or 2 standard fonts (e.g., Arial, Calibri, or Aptos).")
                        else: 
                            st.success("✔ Font family consistency is good.")
                    else: 
                        st.info("Using standard document default fonts.")
                with c2:
                    st.write("**Detected Font Sizes:**")
                    if detected_sizes:
                        for s in sorted(detected_sizes): st.code(s)
                        if len(detected_sizes) > 4: 
                            st.warning("⚠️ Too many different font sizes detected. Consolidate into standard heading and body hierarchy.")
                        else: 
                            st.success("✔ Font sizing hierarchy is well-structured.")
                    else: 
                        st.info("Using standard document default sizes.")
            elif file_type == "PDF":
                st.info("ℹ️ Font extraction is limited on compiled PDFs. Upload your original Word (.docx) or PowerPoint (.pptx) file to audit fonts and sizes!")
            else:
                st.info("ℹ️ Font auditing applies to structured files like Word (.docx) and PowerPoint (.pptx).")

        # --- TAB 3: READABILITY & WORKFLOW STATS ---
        with tab3:
            words = textstat.lexicon_count(document_text, removepunct=True)
            sentences = textstat.sentence_count(document_text)
            grade = textstat.flesch_kincaid_grade(document_text)
            reading_ease = textstat.flesch_reading_ease(document_text)
            
            c1, c2, c3, c4 = st.columns(4)
            c1.metric("Total Word Count", words)
            c2.metric("Total Sentences", sentences)
            c3.metric("Complexity Level", f"Grade {grade}")
            c4.metric("Reading Ease Score", f"{reading_ease} / 100", delta="Formal / Technical" if reading_ease < 50 else "Standard Business")